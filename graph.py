import matplotlib.pyplot as plt
import weather
from pptx.util import Inches
import file


def graph():
    plt.plot()
    x = (weather.locations_name[0:])  # giving data to plot graph
    y = (weather.temperature_list[0:])
    plt.xlabel('Locations')  # taking the x-axis as locations
    plt.ylabel('Temperature')  # taking the y-axis as temperature
    plt.title('Weather report of cities')  # title of the temperature
    plt.plot(x, y, marker="o")
    if len(weather.locations_name) > 5:  # if location_names>5 x_axis name variables will in 60 degree
        plt.xticks(rotation=45)
    # plt.show() the graph
    plt.savefig('graph.jpg')
    title_slide_layout = file.prs.slide_layouts[1]
    slide = file.prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    title.text = "Graph"
    img = 'graph.jpg'
    pic = slide.shapes.add_picture(img, Inches(0), Inches(0), height=Inches(8))
