from pptx import Presentation

prs = Presentation()  # initialising


def ppt_file():  # ppt_file() function is defined under this function slides of ppt & variables used for storing the
    # whether data are defined
    from weather import weather
    var = weather()  # initialising the weather()
    if var is not None:  # to avoid none type errors

        date = var['Date']  # saving the whether data into variables
        time = var['Time']
        Location_name = var['Location_name']
        # temp = var['Temperature']
        wind_speed = var['Wind speed']
        Sky_Description = var['Sky description']
        # weather = var['Weather']
        a=var['moon face']
        temperature= var['temp']
        wind_dir =var['wind direction']


    else:
        return
    title_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(title_slide_layout)  # adding the slide according to count
    title = slide.shapes.title
    title.text = "Weather Forecasting"  # title of the ppt
    subtitle = slide.placeholders[1]
    subtitle.text = (f'Date: {date} \n Time : {time}\n Location Name: {Location_name}\n Temperature : {temperature}⁰C'
                     f'\n Wind speed: {wind_speed}\n wind direction:{wind_dir}'f'\n Sky description: {Sky_Description}\n{a}')
